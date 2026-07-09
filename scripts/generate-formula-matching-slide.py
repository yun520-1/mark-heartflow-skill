#!/usr/bin/env python3
"""Generate a native-PPTX 'formula matching matrix' slide for HeartFlow.

Reads src/formula/formula-triggers.json and builds a native table
(signal class | synonyms/keywords | matched formula refs). No images,
no vector layers -> avoids the old 'Vector zero-size' render error.
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

SKILL_DIR = Path(__file__).resolve().parent.parent
TRIGGERS = SKILL_DIR / "src" / "formula" / "formula-triggers.json"
DOCS = SKILL_DIR / "docs"
OUT = DOCS / "formula-matching-slide.pptx"

SLIDE_W_IN = 10.0
TABLE_W_IN = 9.0


def load_signals():
    data = json.loads(TRIGGERS.read_text(encoding="utf-8"))
    return data["signals"]


def ref_label(ref):
    """Human-readable label for a ref entry."""
    name = ref.get("name")
    if name:
        return name
    kind = ref.get("kind", "")
    r = ref.get("ref", "")
    if r:
        return f"{r} ({kind})" if kind else r
    # ref-less entries: describe by kind
    return kind or "(unspecified)"


def build():
    signals = load_signals()
    prs = Presentation()
    # default slide size is 10x7.5 in; ensure width
    prs.slide_width = Inches(SLIDE_W_IN)
    slide_layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "公式匹配信号矩阵（认知信号 → 公式映射）"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True

    rows = [["信号类", "同义词/关键词", "匹配到的公式 ref 列表"]]
    for key, sig in signals.items():
        terms = sig.get("keywords") or sig.get("synonyms") or []
        terms_str = "、".join(terms) if terms else "—"
        refs = sig.get("refs", [])
        refs_str = "\n".join(f"• {ref_label(r)}" for r in refs) if refs else "—"
        rows.append([key, terms_str, refs_str])

    n_rows = len(rows)
    n_cols = 3
    left = Inches((SLIDE_W_IN - TABLE_W_IN) / 2)
    top = Inches(1.1)
    width = Inches(TABLE_W_IN)
    height = Inches(0.3 * n_rows)
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = graphic_frame.table

    # column widths: signal 1.6, keywords 2.9, refs 4.5
    table.columns[0].width = Inches(1.6)
    table.columns[1].width = Inches(2.9)
    table.columns[2].width = Inches(4.5)

    header_fill = RGBColor(0x2E, 0x5B, 0xA8)
    for c in range(n_cols):
        cell = table.cell(0, c)
        cell.text = rows[0][c]
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(12)
        para.font.bold = True
        para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.vertical_anchor = 1  # middle

    for r in range(1, n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.text = rows[r][c]
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(9)
            cell.vertical_anchor = 1
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF4, 0xF8)

    DOCS.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return n_rows - 1  # data rows = signals count


if __name__ == "__main__":
    n = build()
    print(f"OK rows(signals)={n} file={OUT}")
